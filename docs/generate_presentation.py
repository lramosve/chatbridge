from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# ── Color palette ──
DARK = RGBColor(0x1A, 0x1A, 0x2E)
ACCENT = RGBColor(0x3B, 0x82, 0xF6)
ACCENT2 = RGBColor(0x10, 0xB9, 0x81)
ACCENT3 = RGBColor(0xF5, 0x9E, 0x0B)
ACCENT4 = RGBColor(0xEF, 0x44, 0x44)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_BG = RGBColor(0xF8, 0xFA, 0xFC)
GRAY = RGBColor(0x64, 0x74, 0x8B)
LIGHT_GRAY = RGBColor(0xE2, 0xE8, 0xF0)

def set_slide_bg(slide, color):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_box(slide, left, top, width, height, fill_color, border_color=None, border_width=Pt(1)):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if border_color:
        shape.line.color.rgb = border_color
        shape.line.width = border_width
    else:
        shape.line.fill.background()
    shape.adjustments[0] = 0.05
    return shape

def set_text(shape, text, size=14, bold=False, color=DARK, align=PP_ALIGN.LEFT):
    tf = shape.text_frame
    tf.word_wrap = True
    tf.auto_size = None
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.alignment = align
    return tf

def add_para(tf, text, size=14, bold=False, color=DARK, align=PP_ALIGN.LEFT, space=Pt(4)):
    p = tf.add_paragraph()
    p.text = text
    p.font.size = Pt(size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.alignment = align
    p.space_before = space
    return p

def slide_title(slide, text):
    tb = slide.shapes.add_textbox(Inches(0.8), Inches(0.4), Inches(11), Inches(0.8))
    set_text(tb, text, size=36, bold=True, color=DARK)
    add_box(slide, Inches(0.8), Inches(1.15), Inches(4), Inches(0.05), ACCENT)

def titled_box(slide, left, top, w, h, title, bullets, bg, title_color, bullet_color=DARK):
    box = add_box(slide, left, top, w, h, bg)
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = Pt(14)
    tf.margin_right = Pt(14)
    tf.margin_top = Pt(12)
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(18)
    p.font.bold = True
    p.font.color.rgb = title_color
    for b in bullets:
        bp = tf.add_paragraph()
        bp.text = f"  {b}"
        bp.font.size = Pt(13)
        bp.font.color.rgb = bullet_color
        bp.space_before = Pt(4)
    return box

# ══════════════════════════════════════════════════════════════
# SLIDE 1: Title
# ══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, DARK)
add_box(slide, Inches(0), Inches(3.2), Inches(13.333), Inches(0.06), ACCENT)

tb = slide.shapes.add_textbox(Inches(1), Inches(1.5), Inches(11), Inches(1.5))
tf = set_text(tb, "ChatBridge", size=54, bold=True, color=WHITE)
add_para(tf, "An AI Chat Platform Where Third-Party Apps Live Inside the Conversation", size=22, color=LIGHT_GRAY, space=Pt(14))

tb2 = slide.shapes.add_textbox(Inches(1), Inches(3.8), Inches(11), Inches(2.5))
tf2 = set_text(tb2, "Architecture Deep Dive", size=22, color=ACCENT)
add_para(tf2, "The core problem  •  Key design decisions  •  Trade-offs & uncertainties", size=16, color=GRAY, space=Pt(10))

# ══════════════════════════════════════════════════════════════
# SLIDE 2: The Core Problem
# ══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, LIGHT_BG)
slide_title(slide, "The Core Problem")

# The problem statement - big and clear
box = add_box(slide, Inches(0.8), Inches(1.6), Inches(11.7), Inches(1.6), WHITE, border_color=ACCENT, border_width=Pt(2))
tf = box.text_frame
tf.word_wrap = True
tf.margin_left = Pt(20)
tf.margin_top = Pt(14)
p = tf.paragraphs[0]
p.text = "How do you let unknown third-party apps run inside a chat — rendering UI, calling tools,"
p.font.size = Pt(18)
p.font.color.rgb = DARK
add_para(tf, "maintaining state — while the chatbot stays aware of everything and students stay safe?", size=18, color=DARK, space=Pt(4))

# Why it's hard - three specific challenges with depth
challenges = [
    ("The Boundary Problem", ACCENT4,
     "Third parties control their own UI and tool definitions. The platform cannot predict what any given app will build. "
     "Yet it must render the app, pass it data, track its state, know when it's done, and keep the conversation coherent. "
     "Every handoff — invocation, render, interaction, completion — is a potential failure point."),
    ("The Safety Constraint", ACCENT3,
     "These are children. A broken or malicious app could expose student data, display harmful content, or crash the chat. "
     "The security boundary must be enforceable by the browser, not dependent on trusting app developers. "
     "This rules out any architecture where third-party code runs in the same JavaScript context as the platform."),
    ("The State Gap", ACCENT,
     "The LLM needs to reason about what's happening inside an app it doesn't control. A student asks 'what should I do here?' "
     "during a chess game — the chatbot must know the current board state. But the app owns that state, not the platform. "
     "You need a protocol for the app to share state summaries without the platform becoming coupled to app internals."),
]

y = Inches(3.6)
for title, color, desc in challenges:
    box = add_box(slide, Inches(0.8), y, Inches(11.7), Inches(1.15), WHITE, border_color=color, border_width=Pt(2))
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = Pt(14)
    tf.margin_top = Pt(8)
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(16)
    p.font.bold = True
    p.font.color.rgb = color
    add_para(tf, desc, size=12, color=GRAY, space=Pt(4))
    y += Inches(1.3)

# ══════════════════════════════════════════════════════════════
# SLIDE 3: Decision 1 — Iframe Isolation
# ══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, LIGHT_BG)
slide_title(slide, "Decision 1: Iframe Sandboxing (Not Web Components)")

# The decision
box = add_box(slide, Inches(0.8), Inches(1.6), Inches(11.7), Inches(0.8), RGBColor(0xEF, 0xFA, 0xF5), border_color=ACCENT2, border_width=Pt(2))
tf = box.text_frame
tf.word_wrap = True
tf.margin_left = Pt(16)
tf.margin_top = Pt(8)
p = tf.paragraphs[0]
p.text = 'Each third-party app runs in an <iframe sandbox="allow-scripts allow-forms allow-popups">  —  critically, without allow-same-origin.'
p.font.size = Pt(16)
p.font.bold = True
p.font.color.rgb = ACCENT2

# Three options compared
options = [
    ("Web Components", "Rejected", ACCENT4, [
        "Runs in the same JS context as the platform",
        "A broken component can crash the host page",
        "No true isolation — relies on developer discipline",
        "Unacceptable for third-party code targeting children",
    ]),
    ("Server-Side Rendering", "Rejected", ACCENT3, [
        "App UI rendered on server, streamed as HTML",
        "No client-side interactivity without hydration",
        "Hydration brings us back to needing client-side code",
        "Poor fit for interactive apps like chess",
    ]),
    ("Iframe Sandbox", "Chosen", ACCENT2, [
        "Separate browsing context — true isolation",
        "Browser enforces security (not our code)",
        "Works with any frontend framework",
        "Trade-off: async communication only (postMessage)",
    ]),
]

x = Inches(0.8)
for title, verdict, color, bullets in options:
    box = add_box(slide, x, Inches(2.8), Inches(3.8), Inches(3.0), WHITE, border_color=color, border_width=Pt(2))
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = Pt(14)
    tf.margin_top = Pt(10)
    p = tf.paragraphs[0]
    p.text = f"{title}  —  {verdict}"
    p.font.size = Pt(17)
    p.font.bold = True
    p.font.color.rgb = color
    for b in bullets:
        add_para(tf, f"• {b}", size=12, color=GRAY, space=Pt(5))
    x += Inches(4.0)

# The trade-off callout
box = add_box(slide, Inches(0.8), Inches(6.1), Inches(11.7), Inches(0.9), RGBColor(0xFE, 0xF9, 0xC3), border_color=ACCENT3)
tf = box.text_frame
tf.word_wrap = True
tf.margin_left = Pt(16)
tf.margin_top = Pt(8)
p = tf.paragraphs[0]
p.text = "Trade-off accepted:"
p.font.size = Pt(14)
p.font.bold = True
p.font.color.rgb = ACCENT3
add_para(tf, "Iframes impose friction — apps can't share styling with the parent, communication is async, and apps needing persistent state must "
         "use postMessage to request storage. We accept this friction because the security guarantee is non-negotiable for K-12.", size=12, color=GRAY, space=Pt(3))

# ══════════════════════════════════════════════════════════════
# SLIDE 4: Decision 2 — The PostMessage Protocol
# ══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, LIGHT_BG)
slide_title(slide, "Decision 2: Typed PostMessage Protocol")

# Why this matters
box = add_box(slide, Inches(0.8), Inches(1.6), Inches(11.7), Inches(0.8), RGBColor(0xEB, 0xF5, 0xFF), border_color=ACCENT, border_width=Pt(2))
tf = box.text_frame
tf.word_wrap = True
tf.margin_left = Pt(16)
tf.margin_top = Pt(8)
p = tf.paragraphs[0]
p.text = "This protocol IS the platform. It's the only way apps and the chat can talk. Get it wrong, and nothing works."
p.font.size = Pt(16)
p.font.bold = True
p.font.color.rgb = ACCENT

# Left: the protocol
box = add_box(slide, Inches(0.8), Inches(2.8), Inches(5.6), Inches(4.2), WHITE, border_color=LIGHT_GRAY)
tf = box.text_frame
tf.word_wrap = True
tf.margin_left = Pt(14)
tf.margin_top = Pt(10)
p = tf.paragraphs[0]
p.text = "Seven Message Types"
p.font.size = Pt(18)
p.font.bold = True
p.font.color.rgb = DARK

msgs = [
    ("Platform → App:", True, ACCENT),
    ("  TOOL_INVOKE — call a tool with parameters", False, GRAY),
    ("  RESTORE_STATE — resume after iframe reload", False, GRAY),
    ("", False, GRAY),
    ("App → Platform:", True, ACCENT2),
    ("  APP_READY — iframe loaded, accepting commands", False, GRAY),
    ("  TOOL_RESULT — response data + status", False, GRAY),
    ("  STATE_UPDATE — lightweight state for LLM context", False, GRAY),
    ("  APP_COMPLETE — interaction finished", False, GRAY),
    ("  APP_ERROR — error with details", False, GRAY),
    ("", False, GRAY),
    ("Every message: correlation ID + timestamp + seq #", True, DARK),
    ("Origin validated on both sides, always.", True, DARK),
]

for text, bold, color in msgs:
    add_para(tf, text, size=13, bold=bold, color=color, space=Pt(3))

# Right: the design decisions behind it
box = add_box(slide, Inches(6.9), Inches(2.8), Inches(5.6), Inches(4.2), WHITE, border_color=LIGHT_GRAY)
tf = box.text_frame
tf.word_wrap = True
tf.margin_left = Pt(14)
tf.margin_top = Pt(10)
p = tf.paragraphs[0]
p.text = "Why These Specific Messages?"
p.font.size = Pt(18)
p.font.bold = True
p.font.color.rgb = DARK

decisions = [
    ("STATE_UPDATE is separate from TOOL_RESULT",
     "Apps have ongoing state (chess board) that changes outside of tool calls. The LLM needs this for mid-interaction questions like 'what should I do here?'"),
    ("APP_COMPLETE is explicit, not inferred",
     "Completion signaling is where most teams struggle. We don't try to guess when an app is done — the app tells us. Ambiguity here breaks the conversation flow."),
    ("APP_READY exists because iframes load async",
     "Without it, we'd send TOOL_INVOKE before the app is listening. Race condition on every cold start. APP_READY is the handshake."),
    ("No platform-to-app state push",
     "The app owns its state. The platform only stores a summary. Pushing state creates a sync problem. The app is the source of truth, always."),
]

for title, reason in decisions:
    add_para(tf, title, size=13, bold=True, color=DARK, space=Pt(8))
    add_para(tf, reason, size=11, color=GRAY, space=Pt(2))

# ══════════════════════════════════════════════════════════════
# SLIDE 5: Decision 3 — Extending the AI SDK
# ══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, LIGHT_BG)
slide_title(slide, "Decision 3: Plugin Tools as AI SDK Function Calls")

# The insight
box = add_box(slide, Inches(0.8), Inches(1.6), Inches(11.7), Inches(0.9), RGBColor(0xEF, 0xFA, 0xF5), border_color=ACCENT2, border_width=Pt(2))
tf = box.text_frame
tf.word_wrap = True
tf.margin_left = Pt(16)
tf.margin_top = Pt(8)
p = tf.paragraphs[0]
p.text = "Key insight: Chatbox already uses the Vercel AI SDK with function calling + MCP tools."
p.font.size = Pt(16)
p.font.bold = True
p.font.color.rgb = ACCENT2
add_para(tf, "Plugin tools are just more tools. The LLM doesn't need to know the difference between a built-in web search and a chess move.", size=14, color=GRAY, space=Pt(4))

# The flow diagram as boxes
flow_steps = [
    ("App Manifest", "JSON at /chatbridge-manifest.json\ndeclares tools + parameters\nas JSON Schema", ACCENT, Inches(0.8)),
    ("→", None, GRAY, None),
    ("Plugin Registry", "Fetches + validates manifests\nConverts tool schemas to\nAI SDK tool() format", ACCENT2, Inches(3.65)),
    ("→", None, GRAY, None),
    ("streamText()", "All tools merged:\nbuilt-in + MCP + plugin\nLLM sees one flat list", ACCENT3, Inches(6.5)),
    ("→", None, GRAY, None),
    ("Tool Execute", "Plugin tool's execute()\nsends postMessage to iframe\nawaits TOOL_RESULT", ACCENT4, Inches(9.35)),
]

for item in flow_steps:
    if item[1] is None:  # arrow
        continue
    title, desc, color, x = item
    box = add_box(slide, x, Inches(3.0), Inches(2.6), Inches(1.8), WHITE, border_color=color, border_width=Pt(2))
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = Pt(10)
    tf.margin_top = Pt(8)
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(15)
    p.font.bold = True
    p.font.color.rgb = color
    p.alignment = PP_ALIGN.CENTER
    add_para(tf, desc, size=11, color=GRAY, align=PP_ALIGN.CENTER, space=Pt(6))

# Arrows between boxes
for x_pos in [Inches(3.4), Inches(6.25), Inches(9.1)]:
    a = slide.shapes.add_textbox(x_pos, Inches(3.6), Inches(0.3), Inches(0.5))
    set_text(a, "→", size=24, bold=True, color=GRAY, align=PP_ALIGN.CENTER)

# Why this matters
box = add_box(slide, Inches(0.8), Inches(5.2), Inches(5.6), Inches(1.8), WHITE, border_color=LIGHT_GRAY)
tf = box.text_frame
tf.word_wrap = True
tf.margin_left = Pt(14)
tf.margin_top = Pt(10)
p = tf.paragraphs[0]
p.text = "What we get for free"
p.font.size = Pt(16)
p.font.bold = True
p.font.color.rgb = ACCENT2
for b in ["LLM routes user intent to the right app automatically",
          "Streaming responses pause/resume around tool calls",
          "Multi-provider support (Claude, GPT-4o, etc.)",
          "Context compaction for long conversations"]:
    add_para(tf, f"• {b}", size=12, color=GRAY, space=Pt(4))

box = add_box(slide, Inches(6.9), Inches(5.2), Inches(5.6), Inches(1.8), WHITE, border_color=LIGHT_GRAY)
tf = box.text_frame
tf.word_wrap = True
tf.margin_left = Pt(14)
tf.margin_top = Pt(10)
p = tf.paragraphs[0]
p.text = "Uncertainty: LLM routing accuracy"
p.font.size = Pt(16)
p.font.bold = True
p.font.color.rgb = ACCENT4
for b in ["Ambiguous queries (\"help me practice\") could match multiple apps",
          "Tool descriptions are the main lever — must be precise",
          "May need a routing heuristic or confirmation step",
          "Testing with adversarial prompts is critical"]:
    add_para(tf, f"• {b}", size=12, color=GRAY, space=Pt(4))

# ══════════════════════════════════════════════════════════════
# SLIDE 6: How It Connects — Chess Walkthrough
# ══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, LIGHT_BG)
slide_title(slide, "How It All Connects: Chess Walkthrough")

steps = [
    ("1", '"Let\'s play chess"', "User message", ACCENT),
    ("2", "LLM → chess_start_game", "Function calling routes\nto chess tool", ACCENT3),
    ("3", "Iframe loads chess app", "Sandbox enforced\nAPP_READY received", ACCENT2),
    ("4", "Board renders in chat", "User sees + interacts\nwith chess board", ACCENT),
    ("5", '"What should I do?"', "User asks mid-game", ACCENT3),
]

steps2 = [
    ("6", "App sends STATE_UPDATE", "FEN: rnbqkbnr/pp...\nInjected into LLM context", ACCENT2),
    ("7", "LLM analyzes position", "Responds: \"Consider Nf3\nto control the center\"", ACCENT),
    ("8", "Game ends naturally", "App sends APP_COMPLETE\nwith game result", ACCENT3),
    ("9", "Chat resumes", "LLM has full context\nabout the game played", ACCENT2),
    ("10", "User asks follow-up", "\"Why did I lose?\" — LLM\nreasons from game history", ACCENT),
]

x = Inches(0.4)
box_w = Inches(2.35)
for num, title, desc, color in steps:
    box = add_box(slide, x, Inches(1.5), box_w, Inches(2.2), WHITE, border_color=color, border_width=Pt(2))
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = Pt(10)
    tf.margin_top = Pt(8)
    p = tf.paragraphs[0]
    p.text = num
    p.font.size = Pt(22)
    p.font.bold = True
    p.font.color.rgb = color
    p.alignment = PP_ALIGN.CENTER
    add_para(tf, title, size=14, bold=True, color=DARK, align=PP_ALIGN.CENTER, space=Pt(6))
    add_para(tf, desc, size=11, color=GRAY, align=PP_ALIGN.CENTER, space=Pt(4))
    x += box_w + Inches(0.15)

x = Inches(0.4)
for num, title, desc, color in steps2:
    box = add_box(slide, x, Inches(4.2), box_w, Inches(2.2), WHITE, border_color=color, border_width=Pt(2))
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = Pt(10)
    tf.margin_top = Pt(8)
    p = tf.paragraphs[0]
    p.text = num
    p.font.size = Pt(22)
    p.font.bold = True
    p.font.color.rgb = color
    p.alignment = PP_ALIGN.CENTER
    add_para(tf, title, size=14, bold=True, color=DARK, align=PP_ALIGN.CENTER, space=Pt(6))
    add_para(tf, desc, size=11, color=GRAY, align=PP_ALIGN.CENTER, space=Pt(4))
    x += box_w + Inches(0.15)

# Key callout at bottom
box = add_box(slide, Inches(0.8), Inches(6.6), Inches(11.7), Inches(0.6), RGBColor(0xEB, 0xF5, 0xFF), border_color=ACCENT)
tf = box.text_frame
tf.margin_left = Pt(16)
tf.margin_top = Pt(6)
p = tf.paragraphs[0]
p.text = "Every step uses the same protocol. Weather and Spotify follow the same lifecycle — only the tool schemas and iframe content differ."
p.font.size = Pt(14)
p.font.bold = True
p.font.color.rgb = ACCENT

# ══════════════════════════════════════════════════════════════
# SLIDE 7: Trade-offs & Uncertainties
# ══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, LIGHT_BG)
slide_title(slide, "Trade-offs & Open Questions")

# Trade-offs
tradeoffs = [
    ("Extend Chatbox vs. build from scratch",
     "We chose to fork Chatbox — it gives us streaming, multi-provider LLM, conversation history, and a polished UI for free. "
     "The trade-off: Chatbox is an Electron desktop app. Its architecture (local-first, IndexedDB storage) isn't designed for a "
     "multi-user web platform. We use the web build target and add Supabase for server-side concerns, but some patterns won't "
     "translate cleanly.",
     ACCENT2),
    ("Strict iframe sandbox vs. developer convenience",
     "Omitting allow-same-origin means apps can't use their own cookies or localStorage. This is the right call for safety, but it "
     "means every app that needs persistent state must either use its own backend or ask the platform to store data via postMessage. "
     "That's friction for app developers. We mitigate with an SDK that handles the boilerplate.",
     ACCENT3),
    ("Explicit APP_COMPLETE vs. inferred completion",
     "We require apps to explicitly signal when they're done. An alternative is having the LLM infer completion from context, but that's "
     "unreliable — is a paused chess game 'done'? Explicit signaling is more work for app developers but eliminates a class of bugs "
     "where the chatbot awkwardly resumes too early or waits forever.",
     ACCENT),
]

y = Inches(1.5)
for title, desc, color in tradeoffs:
    box = add_box(slide, Inches(0.8), y, Inches(11.7), Inches(1.3), WHITE, border_color=color, border_width=Pt(2))
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = Pt(14)
    tf.margin_top = Pt(8)
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(16)
    p.font.bold = True
    p.font.color.rgb = color
    add_para(tf, desc, size=12, color=GRAY, space=Pt(4))
    y += Inches(1.45)

# Open questions
box = add_box(slide, Inches(0.8), y, Inches(11.7), Inches(2.1), RGBColor(0xFE, 0xF2, 0xF2), border_color=ACCENT4, border_width=Pt(2))
tf = box.text_frame
tf.word_wrap = True
tf.margin_left = Pt(14)
tf.margin_top = Pt(10)
p = tf.paragraphs[0]
p.text = "Open Questions We Haven't Fully Solved"
p.font.size = Pt(18)
p.font.bold = True
p.font.color.rgb = ACCENT4

questions = [
    "OAuth in popup context — popup blockers vary across browsers. Fallback UX is unclear. Need real-device testing.",
    "Context window pressure at scale — 3 apps with 3 tools each is fine. 20+ apps? We'll need per-conversation tool pruning, which adds complexity.",
    "App review pipeline — MVP uses admin-approved registration. Production K-12 needs automated content scanning. Out of scope for the sprint, but architecturally important.",
    "State restoration reliability — if an app crashes mid-interaction, can it resume from the last STATE_UPDATE? Depends entirely on the app. The platform can't guarantee it.",
]

for q in questions:
    add_para(tf, f"• {q}", size=12, color=GRAY, space=Pt(5))

# ══════════════════════════════════════════════════════════════
# SLIDE 8: Closing
# ══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, DARK)
add_box(slide, Inches(0), Inches(2.8), Inches(13.333), Inches(0.06), ACCENT)

tb = slide.shapes.add_textbox(Inches(1), Inches(0.8), Inches(11), Inches(1.8))
tf = set_text(tb, "Three decisions define this architecture:", size=32, bold=True, color=WHITE)

decisions_summary = [
    ("1. Iframes for isolation", "— safety is enforced by the browser, not by trust"),
    ("2. Typed postMessage protocol", "— the universal contract between platform and apps"),
    ("3. Plugin tools as AI SDK function calls", "— the LLM routes naturally, we don't rebuild the wheel"),
]

y = Inches(3.4)
for title, desc in decisions_summary:
    tb = slide.shapes.add_textbox(Inches(1.2), y, Inches(10.8), Inches(0.6))
    tf = set_text(tb, title, size=22, bold=True, color=ACCENT)
    add_para(tf, desc, size=18, color=LIGHT_GRAY, space=Pt(2))
    y += Inches(0.8)

tb = slide.shapes.add_textbox(Inches(1.2), Inches(6.2), Inches(10.8), Inches(0.6))
tf = set_text(tb, '"A simple chat with one rock-solid integration beats a flashy platform where apps break mid-conversation."',
              size=16, bold=False, color=GRAY, align=PP_ALIGN.LEFT)

# ── Speaker Notes ──
notes = {
    0: (
        "ChatBridge is an AI chat platform where third-party apps live inside the conversation. "
        "A student can play chess, check weather, or create a Spotify playlist — all without leaving the chat window, "
        "with the chatbot staying aware the entire time. Let me walk through the architecture, focusing on the three "
        "design decisions that matter most."
    ),
    1: (
        "The core problem is the boundary between the chat and the third-party app. The platform has to render an app's UI, "
        "call its tools with the right parameters, track its state, know when it's done, and keep the conversation coherent — "
        "all without being able to predict what any given app will do. Consider chess: a student says \"let's play,\" a board "
        "appears, they ask mid-game \"what should I do here?\", the chatbot analyzes the board, the game ends, the conversation "
        "continues. Every handoff in that sequence is a potential failure point.\n\n"
        "This creates three specific hard problems. First, the boundary problem — third parties control their own UI and tools. "
        "The platform must be generic enough to support anything from a calculator to a chess game.\n\n"
        "Second, the safety constraint. These are K-12 students. A broken or malicious app cannot be allowed to access student "
        "data or crash the chat. The security boundary must be enforced by the browser, not by trusting app developers.\n\n"
        "Third, the state gap. The LLM needs to reason about what's happening inside an app it doesn't control. When a student "
        "asks for help mid-chess-game, the chatbot needs to know the board state. But the app owns that state. We need a protocol "
        "for sharing it."
    ),
    2: (
        "The first and most important decision: every third-party app runs in an iframe with the sandbox attribute — and "
        "critically, without allow-same-origin.\n\n"
        "We evaluated three options. Web Components are lighter weight and share styling, but they run in the same JavaScript "
        "context as the platform. A broken component can crash the host page. For third-party code targeting children, that's "
        "unacceptable. Server-side rendering avoids client-side execution, but interactive apps like chess need client-side code, "
        "so it brings us right back to the same problem.\n\n"
        "Iframes give us a separate browsing context. The browser itself enforces the security boundary. Without "
        "allow-same-origin, even malicious JavaScript in the app cannot access the parent window's cookies, localStorage, or DOM. "
        "This is enforced by the browser's Same-Origin Policy, not by our code.\n\n"
        "The trade-off we accepted: this adds friction. Apps can't share styling with the parent. They can't use their own "
        "cookies. Any app that needs persistent state must either use its own backend or ask the platform to store data via "
        "postMessage. We mitigate this with a JavaScript SDK that handles the boilerplate, but the friction is real. We accept "
        "it because the safety guarantee is non-negotiable for K-12."
    ),
    3: (
        "The second decision is the typed postMessage protocol. This is arguably the most important thing we build, because "
        "it IS the platform. It's the only way apps and the chat can communicate.\n\n"
        "Seven message types. Two from platform to app: TOOL_INVOKE sends a tool call with parameters, RESTORE_STATE helps "
        "apps recover after an iframe reload. Five from app to platform: APP_READY signals the iframe loaded, TOOL_RESULT "
        "returns data, STATE_UPDATE shares ongoing state, APP_COMPLETE signals the interaction is finished, and APP_ERROR "
        "reports problems.\n\n"
        "Let me explain why these specific messages exist. STATE_UPDATE is separate from TOOL_RESULT because apps have ongoing "
        "state that changes outside of tool calls — a chess board changes with every move, not just when the LLM invokes a tool. "
        "The LLM needs this state for mid-interaction questions.\n\n"
        "APP_COMPLETE is explicit rather than inferred. We don't try to guess when an app is done. Completion signaling is where "
        "the project brief says most teams struggle, and inferring it from context is unreliable — is a paused chess game \"done\"? "
        "The app tells us, period.\n\n"
        "APP_READY exists because iframes load asynchronously. Without it, we'd send TOOL_INVOKE before the app is listening — "
        "a race condition on every cold start."
    ),
    4: (
        "The third decision: plugin tools are just more function calls to the LLM. Chatbox already uses the Vercel AI SDK with "
        "function calling and MCP tools. Our key insight is that the LLM doesn't need to know the difference between a built-in "
        "web search and a chess move. They're all tools.\n\n"
        "The flow: an app publishes a manifest at a well-known URL declaring its tools as JSON Schema. Our plugin registry "
        "fetches and validates these manifests, then converts them into the AI SDK's tool format. At generation time, all tools — "
        "built-in, MCP, and plugin — are merged into one flat list. When the LLM decides to call a plugin tool, that tool's "
        "execute function sends a postMessage to the app's iframe and awaits the result.\n\n"
        "What we get for free: the LLM routes user intent automatically, streaming pauses and resumes around tool calls, "
        "multi-provider support, and context compaction for long conversations. All of that already works in Chatbox.\n\n"
        "The uncertainty here is routing accuracy. An ambiguous query like \"help me practice\" could match multiple apps. "
        "Tool descriptions are our main lever, but we may need a confirmation step or heuristic. Testing with adversarial "
        "prompts will be critical."
    ),
    5: (
        "Let me show how all three decisions work together. User says \"let's play chess.\" The LLM routes to the chess tool "
        "via function calling — that's decision three. The iframe loads with sandbox restrictions — decision one. The app signals "
        "APP_READY, the board renders. The user plays, asks \"what should I do?\" The app sends a STATE_UPDATE with the FEN "
        "string — that's decision two. The LLM gets the board state in its context, analyzes the position, responds. When the "
        "game ends, APP_COMPLETE fires, and the conversation resumes with full context.\n\n"
        "The critical point: Weather and Spotify follow the exact same lifecycle. Only the tool schemas and iframe content differ. "
        "The protocol is universal."
    ),
    6: (
        "Let me be honest about what we traded away and what we haven't solved.\n\n"
        "Extending Chatbox gives us a mature chat UI for free, but it's a local-first Electron app. Its architecture — IndexedDB, "
        "no shared backend — isn't designed for a multi-user web platform. We bolt on Supabase for server-side concerns, but some "
        "patterns won't translate cleanly.\n\n"
        "The strict iframe sandbox protects students but adds developer friction. Every app needing storage must go through "
        "postMessage or use its own backend. Our SDK helps, but it's still more work than just using localStorage.\n\n"
        "Open questions we haven't fully solved: OAuth popup flows are fragile across browsers — popup blockers vary and fallback "
        "UX is unclear. Context window pressure with 20+ apps would require per-conversation tool pruning. App review pipelines "
        "for production K-12 need automated content scanning — out of scope for the sprint but architecturally important. And "
        "state restoration after an app crash depends entirely on the app; the platform can't guarantee it."
    ),
    7: (
        "Three decisions define this architecture. Iframes for isolation — safety enforced by the browser. A typed postMessage "
        "protocol — the universal contract. Plugin tools as AI SDK function calls — natural routing without reinventing the wheel. "
        "And as the project brief says: a simple chat with one rock-solid integration beats a flashy platform where apps break "
        "mid-conversation. Thank you."
    ),
}

for idx, note_text in notes.items():
    notes_slide = prs.slides[idx].notes_slide
    notes_slide.notes_text_frame.text = note_text

# Save
prs.save('C:\\Users\\lramo\\Documents\\GauntletAI\\Repo\\chatbridge\\docs\\Architecture_Presentation.pptx')
print('Architecture_Presentation.pptx generated successfully.')
